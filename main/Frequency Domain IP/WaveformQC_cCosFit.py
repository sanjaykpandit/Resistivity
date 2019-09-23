import os
from model.data.WaveformDataList import WaveformData
from model.inversion.CosFit import CosFit
from model.Logger import Logger
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# User Input Code Section
# Folder location recorded waveform data
folder = "Samples/H3C/S11T/Waveforms"


# Main Code Starts from Here
logger = Logger()
logger.setDefaultTitle("WaveformQC")

def get_pse(t, v, etol=20):
    y = v
    cs = CosFit(t, y)
    p = cs.get_parameter()
    e = cs.get_error(p)
    return p, cs.get_series(t, p), e

prs = Presentation()

files = os.listdir(folder)
sorted_files = sorted(files, key=lambda x: int(x.split('.')[0]))

for file in sorted_files:
    if file.endswith('.xlsx'):

        logger.InfoLog("Currently Working on {0}".format(file))

        wf = WaveformData()
        wf.load_data(os.path.join(folder, file))

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        w = Inches(2)
        h = Inches(0.25)

        l = Inches(0)
        t = Inches(0)

        txt = slide.shapes.add_textbox(l, t, w, h)
        txt_frame = txt.text_frame
        txt_frame.text = file
        txt_frame.fit_text()

        t = wf.get_time()
        v = wf.get_v(0)

        # cs = CosFit(t, v)
        # p = cs.get_parameter()
        # vf = cs.get_series(t, p)
        # er = cs.get_error(p)

        p, vf, er = get_pse(t, v)

        plt.figure()
        plt.plot(t, v, '.')
        plt.plot(t, vf)
        plt.xlabel("Time (s)")
        plt.ylabel("Voltage (v)")
        plt.title("a: {0:.2f}, f: {1:.2f}, p: {2:.2f}, d: {3:.2f}, e: {4:.2f}".format(p[0], p[1], p[2], p[3], er))
        plt.savefig("__Temp.png")
        plt.close()

        w = Inches(4.67)
        h = Inches(3.5)

        l0 = Inches(0.25)
        t0 = Inches(0.2)
        pic0 = slide.shapes.add_picture("__Temp.png", l0, t0, width=w, height=h)

        v = wf.get_v(1)

        p, vf, er = get_pse(t, v)

        plt.figure()
        plt.plot(t, v, '.')
        plt.plot(t, vf)
        plt.xlabel("Time (s)")
        plt.ylabel("Voltage (v)")
        plt.title("a: {0:.2f}, f: {1:.2f}, p: {2:.2f}, d: {3:.2f}, e: {4:.2f}".format(p[0], p[1], p[2], p[3], er))
        plt.savefig("__Temp.png")
        plt.close()

        l1 = Inches(5.08)
        t1 = Inches(0.2)
        pic1 = slide.shapes.add_picture("__Temp.png", l1, t1, width=w, height=h)

        v = wf.get_v(2)

        p, vf, er = get_pse(t, v)

        plt.figure()
        plt.plot(t, v, '.')
        plt.plot(t, vf)
        plt.xlabel("Time (s)")
        plt.ylabel("Voltage (v)")
        plt.title("a: {0:.2f}, f: {1:.2f}, p: {2:.2f}, d: {3:.2f}, e: {4:.2f}".format(p[0], p[1], p[2], p[3], er))
        plt.savefig("__Temp.png")
        plt.close()

        l2 = Inches(0.25)
        t2 = Inches(3.8)
        pic2 = slide.shapes.add_picture("__Temp.png", l2, t2, width=w, height=h)

        v = wf.get_v(3)

        p, vf, er = get_pse(t, v)

        plt.figure()
        plt.plot(t, v, '.')
        plt.plot(t, vf)
        plt.xlabel("Time (s)")
        plt.ylabel("Voltage (v)")
        plt.title("a: {0:.2f}, f: {1:.2f}, p: {2:.2f}, d: {3:.2f}, e: {4:.2f}".format(p[0], p[1], p[2], p[3], er))
        plt.savefig("__Temp.png")
        plt.close()

        l3 = Inches(5.08)
        t3 = Inches(3.8)
        pic3 = slide.shapes.add_picture("__Temp.png", l3, t3, width=w, height=h)

prs.save(folder + "_QC" + ".pptx")